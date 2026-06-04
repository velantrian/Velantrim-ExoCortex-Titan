"""
🎯 PPTX Parser v2.0 — PowerPoint презентации (НОВЫЙ)
=====================================================
Поддерживает: PPTX, PPTM, ODP

Каскад: python-pptx → Unstructured
"""

import logging
import time

from .base import FileParser, ParseResult

logger = logging.getLogger("velantrim.parsers.pptx")


def _check_available(module_name: str) -> bool:
    import importlib.util
    return importlib.util.find_spec(module_name) is not None


PYTHON_PPTX_AVAILABLE   = _check_available("pptx")
UNSTRUCTURED_AVAILABLE  = _check_available("unstructured")


class PPTXParser(FileParser):
    file_type = "pptx"

    def supported_formats(self) -> list:
        return [".pptx", ".pptm", ".odp"]

    def parse(self, file_path: str) -> ParseResult:
        early = self._check_file(file_path)
        if early is not None:
            return early

        result = ParseResult(file_path=file_path, file_type="pptx")
        result.file_size_bytes = self._get_file_size(file_path)
        start = time.time()

        try:
            if PYTHON_PPTX_AVAILABLE and not file_path.endswith(".odp"):
                text, structured = self._parse_python_pptx(file_path)
                result.extraction_method = "python-pptx"
            elif UNSTRUCTURED_AVAILABLE:
                text, structured = self._parse_unstructured(file_path)
                result.extraction_method = "Unstructured"
            else:
                result.error = (
                    "Ни один PPTX-парсер не доступен. "
                    "pip install python-pptx"
                )
                return result

            result.extracted_text = text
            result.structured_data = structured
            result.word_count = len(text.split()) if text else 0
            result.page_count = structured.get("slide_count", 0)
            result.provenance = self._build_provenance(file_path, result.extraction_method)

            if text.strip():
                result = self._enrich_with_essence(result)

        except Exception as exc:
            logger.error(f"PPTXParser error: {exc}")
            result.error = str(exc)

        result.parse_time_ms = (time.time() - start) * 1000
        return result

    def _parse_python_pptx(self, file_path: str) -> tuple[str, dict]:
        from pptx import Presentation
        prs = Presentation(file_path)

        slides_info: list[dict] = []
        text_parts: list[str] = []
        for i, slide in enumerate(prs.slides):
            slide_text_parts: list[str] = []
            shapes_info: list[dict] = []

            for shape in slide.shapes:
                shape_type = type(shape).__name__
                if shape.has_text_frame:
                    shape_text = "\n".join(
                        para.text for para in shape.text_frame.paragraphs
                    ).strip()
                    if shape_text:
                        slide_text_parts.append(shape_text)
                        shapes_info.append({
                            "type": shape_type,
                            "text_preview": shape_text[:100],
                        })
                # Заметки докладчика
                if hasattr(shape, "has_notes_slide") and shape.has_notes_slide:
                    notes_text = shape.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        slide_text_parts.append(f"[Заметки]: {notes_text}")

            slide_text = "\n".join(slide_text_parts)
            slides_info.append({
                "slide_number": i + 1,
                "shapes": shapes_info,
                "text_length": len(slide_text),
            })

            text_parts.append(f"=== Слайд {i+1} ===\n{slide_text}")

        full_text = "\n\n".join(text_parts)
        return full_text, {
            "format": "pptx",
            "slides": slides_info,
            "slide_count": len(slides_info),
        }

    def _parse_unstructured(self, file_path: str) -> tuple[str, dict]:
        from unstructured.partition.pptx import partition_pptx
        elements = partition_pptx(file_path)
        text = "\n\n".join(str(el) for el in elements)
        return text, {
            "format": "pptx",
            "elements": [
                {"type": type(el).__name__, "text": str(el)}
                for el in elements
            ],
        }
