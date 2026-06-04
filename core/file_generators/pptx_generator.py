"""
🎯 PPTX Generator v1.0 — Презентации PowerPoint
=================================================
Профессиональные презентации из GenerationSpec.

Маппинг блоков на слайды:
- HeadingBlock(level=1) → новый title-slide
- HeadingBlock(level≥2) → секция / заголовок раздела
- ParagraphBlock после heading → bullet на слайде
- TableBlock → отдельный слайд с таблицей
- ImageBlock → отдельный слайд с картинкой
- FactBlock → отдельный слайд (claim как заголовок, метаданные снизу)

Особенности:
- 16:9 widescreen формат по умолчанию
- Темы применяются к фонам, акцентам, шрифтам
- Заметки докладчика для FactBlock
- Pagination в footer
"""

import logging
import time

from .base import (
    CalloutBlock,
    FactBlock,
    FileGenerator,
    GenerationResult,
    GenerationSpec,
    HeadingBlock,
    ImageBlock,
    ListBlock,
    ParagraphBlock,
    QuoteBlock,
    TableBlock,
    get_theme,
)

logger = logging.getLogger("velantrim.generators.pptx")


def _check_available(module_name: str) -> bool:
    import importlib.util
    return importlib.util.find_spec(module_name) is not None


PYTHON_PPTX_AVAILABLE = _check_available("pptx")


class PPTXGenerator(FileGenerator):
    format_name = "pptx"

    def supported_extensions(self) -> list[str]:
        return [".pptx"]

    def generate(
        self,
        spec: GenerationSpec,
        output_path: str,
    ) -> GenerationResult:
        result = GenerationResult(output_path=output_path, format="pptx")
        start = time.time()

        if not PYTHON_PPTX_AVAILABLE:
            result.error = "python-pptx не установлен. pip install python-pptx"
            return result

        self._ensure_output_dir(output_path)
        theme = get_theme(spec.theme)

        try:
            slide_count = self._render(spec, output_path, theme)
            result.method = "python-pptx"
            result.page_count = slide_count
            result.block_count = len(spec.blocks)
            result.file_size_bytes = self._file_size(output_path)
            result.metadata = self._build_provenance(output_path)
        except Exception as exc:
            logger.error(f"PPTX generation error: {exc}")
            result.error = str(exc)

        result.generation_time_ms = (time.time() - start) * 1000
        return result

    def _render(self, spec: GenerationSpec, output_path: str, theme) -> int:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        # 16:9 widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slide_count = 0

        # ─── Title slide ───
        self._add_title_slide(prs, spec.metadata, theme)
        slide_count += 1

        # ─── Группировка блоков в слайды ───
        # Algorithm: HeadingBlock(1) → новый слайд, далее накапливаем контент.
        # HeadingBlock(2+) тоже новый слайд но как раздел.
        # TableBlock, ImageBlock, FactBlock — каждый отдельный слайд.

        current_slide = None
        current_title = ""
        current_bullets: list[str] = []
        slide_index = 1

        def _flush() -> None:
            """Сохранить накопленные bullets на текущий слайд."""
            nonlocal current_slide
            if current_slide is None and (current_title or current_bullets):
                # Создаём content слайд
                current_slide = self._new_content_slide(
                    prs, current_title, current_bullets, theme, slide_index
                )

        for block in spec.blocks:
            if isinstance(block, HeadingBlock):
                # Flush предыдущий слайд
                if current_title or current_bullets:
                    _flush()
                    slide_count += 1
                current_title = block.text
                current_bullets = []
                current_slide = None

            elif isinstance(block, ParagraphBlock):
                current_bullets.append(block.text)

            elif isinstance(block, ListBlock):
                current_bullets.extend(block.items)

            elif isinstance(block, TableBlock):
                # Flush content слайд, потом таблицу отдельным
                if current_title or current_bullets:
                    _flush()
                    slide_count += 1
                    current_title = ""
                    current_bullets = []
                    current_slide = None
                self._add_table_slide(prs, block, theme)
                slide_count += 1

            elif isinstance(block, ImageBlock):
                if current_title or current_bullets:
                    _flush()
                    slide_count += 1
                    current_title = ""
                    current_bullets = []
                    current_slide = None
                self._add_image_slide(prs, block, theme)
                slide_count += 1

            elif isinstance(block, FactBlock):
                if current_title or current_bullets:
                    _flush()
                    slide_count += 1
                    current_title = ""
                    current_bullets = []
                    current_slide = None
                self._add_fact_slide(prs, block, theme)
                slide_count += 1

            elif isinstance(block, CalloutBlock):
                # Callout идёт как bullet с эмодзи
                emoji_map = {"info": "ℹ️", "warning": "⚠️", "success": "✅", "danger": "🚨"}
                emoji = emoji_map.get(block.callout_type, "ℹ️")
                title_part = f"{block.title}: " if block.title else ""
                current_bullets.append(f"{emoji} {title_part}{block.text}")

            elif isinstance(block, QuoteBlock):
                quote = f'"{block.text}"'
                if block.author:
                    quote += f" — {block.author}"
                current_bullets.append(quote)

        # Финальный flush
        if current_title or current_bullets:
            _flush()
            slide_count += 1

        prs.save(output_path)
        return slide_count

    # ─── Слайды ───────────────────────────────────────────────────────────────

    def _add_title_slide(self, prs, metadata, theme):
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Фон под тему
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(theme.background)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(11.33), Inches(2),
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = metadata.title
        run.font.size = Pt(54)
        run.font.bold = True
        run.font.color.rgb = self._hex_to_rgb(theme.primary)
        run.font.name = theme.font_heading

        # Subtitle (author)
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.5),
            Inches(11.33), Inches(1),
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"{metadata.author}"
        run.font.size = Pt(20)
        run.font.color.rgb = self._hex_to_rgb(theme.text_muted)
        run.font.name = theme.font_body

        return slide

    def _new_content_slide(
        self, prs, title: str, bullets: list[str], theme, slide_index: int,
    ):
        from pptx.util import Inches, Pt

        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Фон
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(theme.background)

        # Title bar — цветная полоса сверху
        from pptx.enum.shapes import MSO_SHAPE
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.33), Inches(1),
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self._hex_to_rgb(theme.primary)
        title_bar.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15),
            Inches(12.33), Inches(0.8),
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title or "Содержание"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = self._hex_to_rgb(theme.background)
        run.font.name = theme.font_heading

        # Bullets
        if bullets:
            content_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(1.5),
                Inches(12), Inches(5.5),
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.level = 0
                run = p.add_run()
                run.text = f"• {bullet}"
                run.font.size = Pt(18)
                run.font.color.rgb = self._hex_to_rgb(theme.text)
                run.font.name = theme.font_body
                p.space_after = Pt(theme.spacing_sm)

        return slide

    def _add_table_slide(self, prs, block: TableBlock, theme) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title bar
        if block.caption:
            title_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                Inches(13.33), Inches(0.8),
            )
            title_bar.fill.solid()
            title_bar.fill.fore_color.rgb = self._hex_to_rgb(theme.primary)
            title_bar.line.fill.background()
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.1), Inches(12.33), Inches(0.7),
            )
            p = title_box.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = block.caption
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = self._hex_to_rgb(theme.background)

        # Таблица
        rows = len(block.rows) + 1
        cols = len(block.headers) if block.headers else (
            len(block.rows[0]) if block.rows else 1
        )

        top_offset = Inches(1.2) if block.caption else Inches(0.5)
        table_shape = slide.shapes.add_table(
            rows, cols,
            Inches(0.5), top_offset,
            Inches(12.33), Inches(5.5),
        )
        tbl = table_shape.table

        # Headers
        for i, header in enumerate(block.headers):
            cell = tbl.cell(0, i)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._hex_to_rgb(theme.primary)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.bold = True
                    run.font.color.rgb = self._hex_to_rgb(theme.background)
                    run.font.size = Pt(14)

        # Rows
        for r, row in enumerate(block.rows):
            for c, val in enumerate(row):
                if c < cols:
                    cell = tbl.cell(r + 1, c)
                    cell.text = str(val) if val is not None else ""

    def _add_image_slide(self, prs, block: ImageBlock, theme) -> None:
        import os

        from pptx.util import Inches

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        if not os.path.exists(block.path):
            return

        try:
            slide.shapes.add_picture(
                block.path,
                Inches(1), Inches(1),
                width=Inches(11.33),
            )
            if block.caption:
                from pptx.util import Pt
                cap_box = slide.shapes.add_textbox(
                    Inches(1), Inches(6.8), Inches(11.33), Inches(0.5),
                )
                from pptx.enum.text import PP_ALIGN
                p = cap_box.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = block.caption
                run.font.size = Pt(12)
                run.font.italic = True
                run.font.color.rgb = self._hex_to_rgb(theme.text_muted)
        except Exception as exc:
            logger.warning(f"Image slide failed: {exc}")

    def _add_fact_slide(self, prs, block: FactBlock, theme) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        state_colors = {
            "Validated": theme.success,
            "Supported": theme.success,
            "Hypothesized": theme.warning,
            "Observed": theme.primary,
            "Contradicted": theme.danger,
            "Collapsed": theme.danger,
            "Deprecated": theme.danger,
            "ImmutableCore": theme.success,
        }
        color = state_colors.get(block.epistemic_state, theme.primary)
        conf_label = self._confidence_to_label(block.confidence)

        # Цветная полоса слева
        side_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.3), Inches(7.5),
        )
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = self._hex_to_rgb(color)
        side_bar.line.fill.background()

        # Заголовок — claim
        claim_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.5),
            Inches(12), Inches(3),
        )
        tf = claim_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = block.claim
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = self._hex_to_rgb(theme.text)

        # Метаданные снизу
        meta_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(5),
            Inches(12), Inches(2),
        )
        tf = meta_box.text_frame
        meta_lines = [
            f"📌 ID: {block.fact_id}",
            f"🔵 Состояние: {block.epistemic_state}",
            f"📊 Уверенность: {block.confidence:.3f} ({conf_label})",
            f"📂 Источник: {block.source}",
        ]
        for i, line in enumerate(meta_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(14)
            run.font.color.rgb = self._hex_to_rgb(theme.text_muted)

        # Заметка докладчика
        notes = slide.notes_slide.notes_text_frame
        notes.text = (
            f"Velantrim Fact:\n"
            f"ID: {block.fact_id}\n"
            f"Confidence: {block.confidence}\n"
            f"State: {block.epistemic_state}\n"
            f"Source: {block.source}"
        )

    @staticmethod
    def _hex_to_rgb(hex_color: str):
        from pptx.dml.color import RGBColor
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16),
        )
